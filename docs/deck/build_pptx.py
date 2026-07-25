"""Build the 5-slide summary deck (deck.pptx).

    pip install python-pptx
    python docs/deck/build_pptx.py

Structure maps 1:1 onto the five required elements:
    1 Problem understanding · 2 Solution approach · 3 Key highlights ·
    4 Human-in-the-loop · 5 Challenges & learnings

Layout is authored here rather than parsed from SLIDES.md — the deck is table- and diagram-heavy,
which a markdown bullet-parser cannot lay out. SLIDES.md stays the narrative source and the speaker
notes below are lifted from it.

No screenshots: the deck carries the argument, the demo video carries the evidence.

Visual language matches the app — "Charcoal Ink": near-black canvas, neutral greys, and colour
reserved exclusively for risk semantics (green = cleared, amber = review, red = escalate).

The build asserts every shape stays inside the 16:9 canvas, so a layout regression fails loudly.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "deck.pptx")

# ── palette (mirrors frontend/styles.css) ───────────────────────────────────
BG      = RGBColor(0x0B, 0x0B, 0x0C)
PANEL   = RGBColor(0x14, 0x14, 0x16)
PANEL2  = RGBColor(0x1A, 0x1A, 0x1D)
LINE    = RGBColor(0x32, 0x32, 0x36)
INK     = RGBColor(0xED, 0xED, 0xF0)
MUTED   = RGBColor(0xA1, 0xA1, 0xAA)
FAINT   = RGBColor(0x71, 0x71, 0x7A)
LOW     = RGBColor(0x22, 0xC5, 0x5E)
MED     = RGBColor(0xF5, 0x9E, 0x0B)
HIGH    = RGBColor(0xEF, 0x44, 0x44)
REVIEW  = RGBColor(0xF9, 0x73, 0x16)
ACCENT  = RGBColor(0xE4, 0xE4, 0xE7)
CYAN    = RGBColor(0x0E, 0xA5, 0xE9)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)

W, H = Inches(13.333), Inches(7.5)          # 16:9
FONT = "Inter"
MONO = "Consolas"


# ── primitives ──────────────────────────────────────────────────────────────
def text(slide, x, y, w, h, runs, *, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, line=1.25, font=FONT):
    """runs: str, or list of paragraphs; each paragraph is str or list of (text, {opts})."""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(runs if isinstance(runs, list) else [runs]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        for t, o in (para if isinstance(para, list) else [(para, {})]):
            r = p.add_run()
            r.text = t
            r.font.name = o.get("font", font)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.color.rgb = o.get("color", color)
    return tf


def rect(slide, x, y, w, h, fill=PANEL, outline=LINE, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if outline is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = outline
        shp.line.width = Pt(0.75)
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.shadow.inherit = False
    return shp


def slide_base(prs, kicker, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    rect(s, 0, 0, W, Pt(3), fill=ACCENT, outline=None, radius=False)
    text(s, Inches(0.6), Inches(0.38), Inches(11), Inches(0.28),
         kicker.upper(), size=11, color=FAINT, bold=True)
    text(s, Inches(0.6), Inches(0.68), Inches(12.2), Inches(0.55),
         title, size=28, color=INK, bold=True)
    if subtitle:
        text(s, Inches(0.6), Inches(1.22), Inches(12.2), Inches(0.4),
             subtitle, size=12.5, color=MUTED)
    return s


def note(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def label(slide, x, y, w, txt):
    text(slide, x, y, w, Inches(0.22), txt.upper(), size=9.5, color=FAINT, bold=True)


def stat(slide, x, y, w, lbl, value, color=INK, delta=None, h=Inches(1.0)):
    rect(slide, x, y, w, h, fill=PANEL)
    text(slide, x + Inches(0.16), y + Inches(0.11), w - Inches(0.3), Inches(0.2),
         lbl.upper(), size=8.5, color=FAINT, bold=True)
    text(slide, x + Inches(0.16), y + Inches(0.32), w - Inches(0.3), Inches(0.45),
         value, size=23, color=color, bold=True)
    if delta:
        text(slide, x + Inches(0.16), y + Inches(0.73), w - Inches(0.3), Inches(0.2),
             delta, size=8.5, color=MUTED)


def flow_box(slide, x, y, w, h, title, lines, *, accent=LINE, tcolor=INK):
    rect(slide, x, y, w, h, fill=PANEL, outline=accent)
    text(slide, x + Inches(0.14), y + Inches(0.11), w - Inches(0.28), Inches(0.24),
         title, size=10.5, color=tcolor, bold=True)
    text(slide, x + Inches(0.14), y + Inches(0.38), w - Inches(0.28), h - Inches(0.48),
         lines, size=8.5, color=MUTED, space_after=2, line=1.15)


def arrow(slide, x, y, w=Inches(0.28)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.16))
    a.fill.solid(); a.fill.fore_color.rgb = LINE
    a.line.fill.background(); a.shadow.inherit = False
    return a


def card(slide, x, y, w, h, head, body, *, accent=None, hsize=10.5, bsize=8.5,
         hcolor=INK, fill=PANEL, bar=False):
    rect(slide, x, y, w, h, fill=fill, outline=accent if (accent and not bar) else LINE)
    if bar and accent:
        rect(slide, x, y, Pt(3), h, fill=accent, outline=None, radius=False)
    pad = Inches(0.2) if bar else Inches(0.15)
    text(slide, x + pad, y + Inches(0.07), w - pad - Inches(0.14), Inches(0.24),
         head, size=hsize, color=hcolor, bold=True)
    text(slide, x + pad, y + Inches(0.29), w - pad - Inches(0.12), h - Inches(0.36),
         body, size=bsize, color=MUTED, line=1.16, space_after=2)


# ── slide 1 · problem understanding ─────────────────────────────────────────
def slide1(prs):
    s = slide_base(prs, "01 · The Problem",
                   "Too many alarms. Almost none of them real.",
                   "The hard part is not spotting alarms — it is sorting the real ones from the noise, and proving afterwards why you decided what you did.")

    label(s, Inches(0.6), Inches(1.72), Inches(6.0), "what the industry looks like today")
    facts = [
        ("9 in 10", "alarms turn out to be nothing", "so the real work is sorting, not spotting", HIGH),
        ("2–4%", "of alarms are worth acting on", "almost all the effort leads nowhere", MED),
        ("70–80%", "of a reviewer's day is spent on dead ends", "the scarce thing is attention, not data", REVIEW),
        ("2–5%", "of the world's money is dirty (UN estimate)", "and missing one is the expensive mistake", MUTED),
    ]
    for i, (n, what, why, c) in enumerate(facts):
        ry = Inches(2.0) + Inches(0.55 * i)
        rect(s, Inches(0.6), ry, Inches(6.0), Inches(0.47), fill=PANEL)
        text(s, Inches(0.74), ry + Inches(0.08), Inches(1.05), Inches(0.3),
             n, size=13, color=c, bold=True)
        text(s, Inches(1.88), ry + Inches(0.05), Inches(4.6), Inches(0.2),
             what, size=9, color=INK)
        text(s, Inches(1.88), ry + Inches(0.24), Inches(4.6), Inches(0.2),
             why, size=8, color=FAINT)
    text(s, Inches(0.6), Inches(4.24), Inches(6.0), Inches(0.36),
         "Rough industry figures, not audited numbers — but every source tells the same story.",
         size=8.5, color=FAINT)

    label(s, Inches(6.9), Inches(1.72), Inches(5.83), "one customer arrives as five separate files")
    rows = [
        ("kyc.txt",          "form-like", "who they say they are, and what they do"),
        ("account.txt",      "form-like", "how long they have banked here, and where"),
        ("transactions.txt", "table",     "what they actually did — hundreds of rows"),
        ("rm_notes.txt",     "prose",     "their banker's gut feeling, written out"),
        ("id_document.txt",  "prose",     "what the ID scan actually said"),
    ]
    for i, (f, kind, hides) in enumerate(rows):
        ry = Inches(2.0) + Inches(0.44 * i)
        rect(s, Inches(6.9), ry, Inches(5.83), Inches(0.38), fill=PANEL)
        text(s, Inches(7.04), ry + Inches(0.07), Inches(1.7), Inches(0.24),
             f, size=9.5, color=INK, bold=True, font=MONO)
        text(s, Inches(8.76), ry + Inches(0.08), Inches(0.92), Inches(0.22),
             kind, size=8, color=REVIEW if kind == "prose" else FAINT, bold=True)
        text(s, Inches(9.72), ry + Inches(0.08), Inches(2.9), Inches(0.22),
             hides, size=8.5, color=MUTED)
    text(s, Inches(6.9), Inches(4.24), Inches(5.83), Inches(0.36),
         [[("The two files that give the game away are plain English. ",
            {"size": 9, "color": MUTED}),
           ("Traditional software cannot read them.", {"size": 9, "color": REVIEW, "bold": True})]])

    label(s, Inches(0.6), Inches(4.76), Inches(12.13), "so three things go wrong")
    fails = [
        ("It is slow", HIGH,
         "Minutes per customer, and nobody knows who to open first — so the worst case might be read last."),
        ("It is inconsistent", MED,
         "Give the same file to two reviewers and you get two different answers. Nobody agrees on the yardstick."),
        ("Nobody can explain it", REVIEW,
         "The thinking stays in someone's head. Ask six months later why a customer was let through, and you get a shrug."),
    ]
    for i, (h, c, b) in enumerate(fails):
        card(s, Inches(0.6) + Inches(4.09) * i, Inches(5.02), Inches(3.93), Inches(0.78),
             h, b, accent=c, bar=True, hcolor=c, hsize=11.5, bsize=8.5)

    cy = Inches(5.98)
    rect(s, Inches(0.6), cy, Inches(12.13), Inches(1.06), fill=PANEL2, outline=ACCENT)
    text(s, Inches(0.84), cy + Inches(0.11), Inches(11.6), Inches(0.22),
         "WHAT WE SET OUT TO BUILD", size=9, color=FAINT, bold=True)
    text(s, Inches(0.84), cy + Inches(0.34), Inches(11.65), Inches(0.66),
         [[("Do the first pass. ", {"bold": True, "color": INK, "size": 11}),
           ("Read every file, dig like a good reviewer would, put the whole customer list in order of "
            "risk, and show the working — so a person starts with evidence instead of a blank page.",
            {"color": MUTED, "size": 11})],
          [("Missing a real bad actor is a disaster; a false alarm is just costly. So the system has "
            "to ", {"color": MUTED, "size": 11}),
           ("admit when it is unsure and hand over", {"bold": True, "color": REVIEW, "size": 11}),
           (", instead of bluffing. That one rule shaped the whole design.",
            {"color": MUTED, "size": 11})]], space_after=4, line=1.25)

    note(s, "Lead with the reviewer's desk, not the tech. The numbers establish this is a sorting and "
            "explaining problem, not a spotting problem. Everything in the next four slides answers "
            "one of the three things that go wrong.")
    return s


# ── slide 2 · solution approach ─────────────────────────────────────────────
def slide2(prs):
    s = slide_base(prs, "02 · How We Built It",
                   "A team of AI specialists, then one lead investigator",
                   "Four steps per customer. No scoring formula anywhere in the code — the AI decides the number; the code hands it the tools, the memory, and rules it cannot break.")

    py = Inches(1.76)
    bw, bh = Inches(2.72), Inches(1.42)
    gap = Inches(0.36)
    xs = [Inches(0.6) + (bw + gap) * i for i in range(4)]
    flow_box(s, xs[0], py, bw, bh, "1 · GATHER THE FILE",
             ["five separate files in,", "one tidy case folder out", "",
              "→ nothing is left unread"],
             accent=CYAN, tcolor=RGBColor(0x7D, 0xD3, 0xFC))
    flow_box(s, xs[1], py, bw, bh, "2 · REMEMBER",
             ["what we knew about them before", "how we handled similar people",
              "lessons from past corrections", "", "the AML rulebook"],
             accent=TEAL, tcolor=RGBColor(0x5E, 0xEA, 0xD4))
    flow_box(s, xs[2], py, bw, bh, "3 · INVESTIGATE",
             ["3 specialists work AT ONCE", "   background · money · papers", "",
              "1 lead works STEP BY STEP", "   asks, reads, then decides"],
             accent=VIOLET, tcolor=RGBColor(0xC4, 0xB5, 0xFD))
    flow_box(s, xs[3], py, bw, bh, "4 · DECIDE + LEARN",
             ["sure enough → decide it", "not sure → ask a person", "",
              "their answer → remembered"], accent=REVIEW, tcolor=RGBColor(0xFD, 0xBA, 0x74))
    for i in range(3):
        arrow(s, xs[i] + bw + Inches(0.04), py + Inches(0.63))

    ty = Inches(3.44)
    label(s, Inches(0.6), ty, Inches(6.0), "why a team, and then a lead")
    text(s, Inches(1.9), ty + Inches(0.26), Inches(2.1), Inches(0.22),
         "The 3 specialists", size=9.5, color=ACCENT, bold=True)
    text(s, Inches(4.05), ty + Inches(0.26), Inches(2.4), Inches(0.22),
         "The lead investigator", size=9.5, color=ACCENT, bold=True)
    trows = [
        ("How they work", "all three at once", "one step at a time"),
        ("What they see", "only their own topic", "all 3 views + the whole file"),
        ("What they give", "an opinion, not a verdict", "the final call, with reasons"),
        ("Why this way", "fast, and no groupthink", "each answer shapes the next question"),
    ]
    for i, (a, b, c) in enumerate(trows):
        ry = ty + Inches(0.52) + Inches(0.31 * i)
        if i % 2 == 0:
            rect(s, Inches(0.6), ry - Inches(0.03), Inches(5.86), Inches(0.29),
                 fill=PANEL, outline=None)
        text(s, Inches(0.72), ry, Inches(1.15), Inches(0.24), a, size=8.5, color=FAINT, bold=True)
        text(s, Inches(1.9), ry, Inches(2.1), Inches(0.24), b, size=8.5, color=MUTED)
        text(s, Inches(4.05), ry, Inches(2.35), Inches(0.24), c, size=8.5, color=MUTED)

    label(s, Inches(6.85), ty, Inches(5.88), "five kinds of memory it can draw on")
    tiers = [
        ("Scratchpad",   "notes it takes while working",  "wiped when done",     REVIEW),
        ("This person",  "how we judged them before",     "kept forever",        FAINT),
        ("Past cases",   "people who looked like this",   "kept forever",        TEAL),
        ("The rulebook", "what each crime pattern means", "fixed reference",     CYAN),
        ("Lessons",      "where we got it wrong before",  "grows as people fix", VIOLET),
    ]
    for i, (name, what, life, col) in enumerate(tiers):
        ry = ty + Inches(0.28) + Inches(0.37 * i)
        rect(s, Inches(6.85), ry, Inches(5.88), Inches(0.33),
             fill=PANEL if i % 2 == 0 else BG, outline=None)
        rect(s, Inches(6.97), ry + Inches(0.11), Inches(0.1), Inches(0.1), fill=col, outline=None)
        text(s, Inches(7.18), ry + Inches(0.05), Inches(1.5), Inches(0.24),
             name, size=9, color=INK, bold=True)
        text(s, Inches(8.85), ry + Inches(0.05), Inches(2.6), Inches(0.24),
             what, size=8.5, color=MUTED)
        text(s, Inches(11.5), ry + Inches(0.05), Inches(1.2), Inches(0.24),
             life, size=8, color=FAINT)

    ky = Inches(5.44)
    rect(s, Inches(0.6), ky, Inches(6.0), Inches(1.06), fill=PANEL2, outline=VIOLET)
    text(s, Inches(0.82), ky + Inches(0.11), Inches(5.6), Inches(0.22),
         "THE DECISION WE ARGUED ABOUT MOST", size=9, color=FAINT, bold=True)
    text(s, Inches(0.82), ky + Inches(0.34), Inches(5.62), Inches(0.66),
         [[("The lead works ", {"size": 9.5, "color": MUTED}),
           ("one step at a time", {"size": 9.5, "color": INK, "bold": True}),
           (" because real investigation is like that — you cannot know which document to open "
            "until you have seen the money. The three specialists work ",
            {"size": 9.5, "color": MUTED}),
           ("all at once", {"size": 9.5, "color": INK, "bold": True}),
           (" because their topics never overlap.", {"size": 9.5, "color": MUTED})]], line=1.2)

    rect(s, Inches(6.85), ky, Inches(5.88), Inches(1.06), fill=PANEL2, outline=TEAL)
    text(s, Inches(7.07), ky + Inches(0.11), Inches(5.5), Inches(0.22),
         "WHY IT REMEMBERS BEFORE IT THINKS", size=9, color=FAINT, bold=True)
    text(s, Inches(7.07), ky + Inches(0.34), Inches(5.5), Inches(0.66),
         [[("Memory is loaded ", {"size": 9.5, "color": MUTED}),
           ("before", {"size": 9.5, "color": INK, "bold": True}),
           (" anyone starts thinking, so every part of the pipeline works from the same background — "
            "past judgements, similar people, and lessons a human taught it. That is the difference "
            "between a system and a clever prompt.", {"size": 9.5, "color": MUTED})]], line=1.2)

    note(s, "The split — a team working at once, then one lead working step by step — is the core "
            "architectural claim. The memory table is the second: what it learns outlives the request.")
    return s


# ── slide 3 · key highlights ────────────────────────────────────────────────
def slide3(prs):
    s = slide_base(prs, "03 · What It Actually Does",
                   "It digs, it points at the evidence, it shows the working",
                   "A real investigation of CUST_018 — a customer whose money did not match his story. Every step below was recorded automatically. That record is the receipt.")

    label(s, Inches(0.6), Inches(1.72), Inches(6.2), "what it did, in its own order — 12 steps")
    steps = [
        ("1",   "read_document",      "opened the ID scan — passport checks out"),
        ("2",   "read_document",      "read the banker's notes — “can't explain his income”"),
        ("3",   "query_transactions", "pulled all 31 payments, £62,174 coming in"),
        ("4-9", "query_transactions", "sliced them six ways — cash, who, when"),
        ("10",  "find_txn_patterns",  "spotted a possible pattern in rows S00–S03"),
        ("11",  "note",               "wrote itself: “4 deposits just under the limit, 6 days”"),
        ("12",  "finalize",           "decided: 83 · HIGH RISK · send to a senior · 85% sure"),
    ]
    for i, (n, tool, what) in enumerate(steps):
        ry = Inches(2.0) + Inches(0.38 * i)
        last = i == len(steps) - 1
        rect(s, Inches(0.6), ry, Inches(6.2), Inches(0.32),
             fill=PANEL2 if last else PANEL, outline=HIGH if last else None)
        text(s, Inches(0.7), ry + Inches(0.05), Inches(0.4), Inches(0.22),
             n, size=8, color=FAINT, bold=True, align=PP_ALIGN.RIGHT)
        text(s, Inches(1.2), ry + Inches(0.05), Inches(1.72), Inches(0.22),
             tool, size=8.5, color=HIGH if last else RGBColor(0xC4, 0xB5, 0xFD),
             bold=True, font=MONO)
        text(s, Inches(2.96), ry + Inches(0.05), Inches(3.76), Inches(0.22),
             what, size=8.5, color=INK if last else MUTED)

    label(s, Inches(7.15), Inches(1.72), Inches(5.58), "the answer comes with its receipts attached")
    oy = Inches(2.0)
    rect(s, Inches(7.15), oy, Inches(5.58), Inches(1.4), fill=PANEL2, outline=HIGH)
    text(s, Inches(7.35), oy + Inches(0.08), Inches(5.2), Inches(0.4),
         [[("83", {"size": 26, "bold": True, "color": HIGH}),
           (" / 100", {"size": 11, "color": FAINT}),
           ("    HIGH RISK · SEND TO A SENIOR", {"size": 12, "bold": True, "color": HIGH}),
           ("   85% sure", {"size": 10, "color": MUTED})]])
    text(s, Inches(7.35), oy + Inches(0.52), Inches(5.22), Inches(0.82),
         [[("Proof it points at:  ", {"size": 8.5, "color": FAINT}),
           ("payments S00, S01, S02, S03  ·  the banker's notes",
            {"size": 8.5, "color": LOW, "font": MONO})],
          [("Why:  ", {"size": 8.5, "color": FAINT}),
           ("home country carries elevated risk · line of work is high-risk · four cash deposits "
            "totalling $37k, each kept just under the reporting limit, all inside six days",
            {"size": 8.5, "color": MUTED})]],
         space_after=3, line=1.25)

    label(s, Inches(7.15), Inches(3.58), Inches(5.58), "four rules it is not allowed to break")
    guards = [
        ("It cannot make up proof", "every reference is checked against what the tools really returned — invented ones are thrown out"),
        ("It cannot ramble forever", "12 steps maximum; run past that and a person picks it up"),
        ("It cannot come back empty", "even if something crashes, you still get a real decision, routed to a person"),
        ("Tools suggest, AI decides", "the pattern scanner only says “this looks odd” — the AI has to agree and point at the rows"),
    ]
    for i, (h, b) in enumerate(guards):
        card(s, Inches(7.15), Inches(3.86) + Inches(0.55 * i), Inches(5.58), Inches(0.5),
             h, b, hsize=9.5, bsize=8, hcolor=LOW)

    ry = Inches(4.86)
    label(s, Inches(0.6), ry, Inches(6.2), "and it ran the whole book, start to finish")
    cw = Inches(1.19)
    kpis = [("customers", "22", INK), ("cleared on its own", "27%", LOW), ("sent up", "3", HIGH),
            ("asked a human", "1", REVIEW), ("tests pass", "14 ✓", LOW)]
    for i, (l, v, c) in enumerate(kpis):
        x = Inches(0.6) + (cw + Inches(0.06)) * i
        rect(s, x, ry + Inches(0.26), cw, Inches(0.62), fill=PANEL)
        text(s, x + Inches(0.1), ry + Inches(0.32), cw - Inches(0.2), Inches(0.2),
             l.upper(), size=7, color=FAINT, bold=True)
        text(s, x + Inches(0.1), ry + Inches(0.5), cw - Inches(0.2), Inches(0.3),
             v, size=15, color=c, bold=True)

    label(s, Inches(0.6), Inches(6.0), Inches(6.2), "same AI, two very different customers")
    rect(s, Inches(0.6), Inches(6.26), Inches(6.2), Inches(0.42), fill=PANEL, outline=HIGH)
    text(s, Inches(0.76), Inches(6.34), Inches(5.9), Inches(0.26),
         [[("CUST_018", {"size": 9.5, "bold": True, "color": INK, "font": MONO}),
           ("  odd cash, odd story → ", {"size": 8.5, "color": MUTED}),
           ("83 · SEND TO A SENIOR", {"size": 9.5, "bold": True, "color": HIGH}),
           (" · 85% sure", {"size": 8.5, "color": MUTED})]])
    rect(s, Inches(0.6), Inches(6.76), Inches(6.2), Inches(0.42), fill=PANEL, outline=LOW)
    text(s, Inches(0.76), Inches(6.84), Inches(5.9), Inches(0.26),
         [[("CUST_000", {"size": 9.5, "bold": True, "color": INK, "font": MONO}),
           ("  teacher, salary in, rent out → ", {"size": 8.5, "color": MUTED}),
           ("5 · CLEARED", {"size": 9.5, "bold": True, "color": LOW}),
           (" · 95% sure", {"size": 8.5, "color": MUTED})]])

    note(s, "Walk the steps out loud. This is 'show your working' made real. Point at step 10 — the "
            "tool only says something looks odd; the AI had to agree and name the exact payments.")
    return s


# ── slide 4 · human-in-the-loop ─────────────────────────────────────────────
def slide4(prs):
    s = slide_base(prs, "04 · Where People Come In",
                   "When it is unsure, it asks. And it remembers your answer.",
                   "The AI says how sure it is. If it is not sure enough, a person decides instead — and that correction makes the whole system better, three different ways.")

    gy = Inches(1.78)
    rect(s, Inches(0.6), gy, Inches(2.4), Inches(0.66), fill=PANEL2, outline=VIOLET)
    text(s, Inches(0.74), gy + Inches(0.09), Inches(2.15), Inches(0.5),
         [[("A decision", {"size": 11, "bold": True, "color": INK})],
          [("+ how sure the AI says it is", {"size": 8, "color": MUTED})]], space_after=1)
    arrow(s, Inches(3.08), gy + Inches(0.25))
    rect(s, Inches(3.5), gy, Inches(2.1), Inches(0.66), fill=PANEL, outline=ACCENT)
    text(s, Inches(3.6), gy + Inches(0.18), Inches(1.9), Inches(0.32),
         "at least 60% sure?", size=10.5, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(5.92), Inches(1.66), Inches(2.66), Inches(0.55), fill=PANEL, outline=LOW)
    text(s, Inches(6.06), Inches(1.77), Inches(2.4), Inches(0.32),
         [[("YES → ", {"size": 9.5, "bold": True, "color": LOW}),
           ("it decides on its own", {"size": 9.5, "color": MUTED})]])
    rect(s, Inches(5.92), Inches(2.31), Inches(2.66), Inches(0.55), fill=PANEL, outline=REVIEW)
    text(s, Inches(6.06), Inches(2.42), Inches(2.4), Inches(0.32),
         [[("NO → ", {"size": 9.5, "bold": True, "color": REVIEW}),
           ("it goes in a person's queue", {"size": 9.5, "color": MUTED})]])
    arrow(s, Inches(8.7), Inches(2.44))
    rect(s, Inches(9.1), Inches(2.22), Inches(3.63), Inches(0.66), fill=PANEL2, outline=REVIEW)
    text(s, Inches(9.26), Inches(2.33), Inches(3.35), Inches(0.5),
         [[("A person sets the right answer", {"size": 10, "bold": True, "color": INK})],
          [("and sees everything the AI saw", {"size": 8, "color": MUTED})]], space_after=1)

    wy = Inches(3.16)
    rect(s, Inches(0.6), wy, Inches(6.0), Inches(1.34), fill=PANEL)
    text(s, Inches(0.78), wy + Inches(0.11), Inches(5.6), Inches(0.22),
         "WHY WE ASK “HOW SURE?” AND NOT “HOW BAD?”", size=9, color=FAINT, bold=True)
    text(s, Inches(0.78), wy + Inches(0.36), Inches(5.64), Inches(0.92),
         [[("A middling score is not automatically a doubtful one — the AI can be very sure a "
            "customer is a plain, boring 58. What matters is whether ", {"size": 9.5, "color": MUTED}),
           ("the evidence really backs the answer", {"size": 9.5, "color": INK, "bold": True}),
           (". We ask the AI to rate its own certainty and we tell it plainly that shaky cases go to "
            "a person — so ", {"size": 9.5, "color": MUTED}),
           ("owning up to doubt is the winning move", {"size": 9.5, "color": LOW, "bold": True}),
           (", not a failure.", {"size": 9.5, "color": MUTED})]], line=1.2)

    ry2 = Inches(4.68)
    rect(s, Inches(0.6), ry2, Inches(6.0), Inches(1.06), fill=PANEL)
    text(s, Inches(0.78), ry2 + Inches(0.11), Inches(5.6), Inches(0.22),
         "THE PERSON SEES EVERYTHING — NOTHING IS HIDDEN", size=9, color=FAINT, bold=True)
    text(s, Inches(0.78), ry2 + Inches(0.35), Inches(5.64), Inches(0.64),
         "The suggested score and how sure it was · all three specialists' views, including where "
         "they disagreed · every step it took · the notes it wrote itself · and exactly what it "
         "remembered before starting.", size=9.5, color=MUTED, line=1.2)

    label(s, Inches(6.85), Inches(3.16), Inches(5.88), "one correction teaches it three ways")
    paths = [
        ("1 · a worked example", "the next customer who looks like this gets shown your answer", TEAL),
        ("2 · a written lesson", "turned into a rule of thumb it reads before every future case", VIOLET),
        ("3 · a note on that person's file", "next time this customer comes round, it sees what changed", CYAN),
    ]
    for i, (h, b, c) in enumerate(paths):
        card(s, Inches(6.85), Inches(3.44) + Inches(0.6 * i), Inches(5.88), Inches(0.54),
             h, b, accent=c, bar=True, hsize=9.5, bsize=8.5)
    text(s, Inches(6.85), Inches(5.24), Inches(5.88), Inches(0.24),
         [[("Most systems manage one of these. Three means your single correction improves ",
            {"size": 8.5, "color": FAINT}),
           ("similar cases, every case, and this one", {"size": 8.5, "color": MUTED, "bold": True}),
           (".", {"size": 8.5, "color": FAINT})]])

    ay = Inches(5.62)
    rect(s, Inches(6.85), ay, Inches(5.88), Inches(1.16), fill=PANEL2, outline=LOW)
    text(s, Inches(7.05), ay + Inches(0.11), Inches(5.6), Inches(0.22),
         "IT ONLY LEARNS FROM PEOPLE, NEVER FROM ITSELF", size=9, color=LOW, bold=True)
    text(s, Inches(7.05), ay + Inches(0.35), Inches(5.6), Inches(0.74),
         "It only copies from cases a human actually checked. If it were allowed to learn from its "
         "own unchecked work, one early mistake would quietly become the house rule. Every decision — "
         "cleared or escalated — is written down permanently and cannot be edited later.",
         size=9.5, color=MUTED, line=1.2)

    note(s, "This is the slide that answers 'does it get better?'. Emphasise the three paths, then "
            "the learn-only-from-people rule — that is the detail showing we saw the trap coming.")
    return s


# ── slide 5 · challenges & learnings ────────────────────────────────────────
def slide5(prs):
    s = slide_base(prs, "05 · What Broke, and What We Learned",
                   "Four things went wrong. Each one made the design better.",
                   "These are not hypotheticals — they are real bugs that shipped, got caught, and got fixed.")

    label(s, Inches(0.6), Inches(1.72), Inches(6.6), "what broke")
    chals = [
        ("1 · The investigator that never wrapped up", HIGH,
         "It ran out of steps mid-investigation and handed back a blank “0, not sure” — which the "
         "system happily read as a genuine answer. A failure wearing the costume of a result.",
         "Fix: in the last two steps, the only thing it is allowed to do is finish. A time limit "
         "needs a forced ending, not just a buzzer."),
        ("2 · The feature we deleted that kept coming back", REVIEW,
         "We cut watchlist screening from the project, yet it kept turning up in the output — hiding "
         "in three places: a leftover default, a stray sentence in a reference file, and an old "
         "cached answer.",
         "Fix: clear the data, not just the code. Then re-ran all 22 customers to prove it was gone."),
        ("3 · “It must be the API throttling us”", MED,
         "Scoring several customers at once took 76–90 seconds and I blamed the provider. The "
         "measurements said otherwise: four calls at once took 1.55s versus 1.48s for one. The real "
         "cause was our own code rebuilding the same tools on every single call.",
         "Fix: build them once and reuse. That step went from 7.5 seconds to 5 milliseconds."),
        ("4 · It was about to start marking its own homework", VIOLET,
         "It would have treated its own unchecked past decisions as precedent — turning one early "
         "mistake into a permanent habit.",
         "Fix: it may only learn from cases a real person signed off."),
    ]
    for i, (h, c, problem, fix) in enumerate(chals):
        cy = Inches(2.0) + Inches(1.26 * i)
        rect(s, Inches(0.6), cy, Inches(6.6), Inches(1.14), fill=PANEL)
        rect(s, Inches(0.6), cy, Pt(3), Inches(1.14), fill=c, outline=None, radius=False)
        text(s, Inches(0.82), cy + Inches(0.08), Inches(6.25), Inches(0.22),
             h, size=10, color=c, bold=True)
        text(s, Inches(0.82), cy + Inches(0.3), Inches(6.28), Inches(0.52),
             problem, size=8.5, color=MUTED, line=1.14)
        text(s, Inches(0.82), cy + Inches(0.86), Inches(6.28), Inches(0.24),
             fix, size=8.5, color=LOW, line=1.14)

    label(s, Inches(7.5), Inches(1.72), Inches(5.23), "what we learned")
    lessons = [
        ("Hand the AI facts, not conclusions",
         "Tools that say “this looks odd, here are the rows” get better thinking than tools that say "
         "“this is fraud” — and you can check the answer afterwards."),
        ("“How sure are you?” beats “how bad is it?”",
         "A confident middling score needs nobody. A shaky low one needs a person. Routing on doubt, "
         "not on severity, is what makes escalation mean something."),
        ("Measure before you fix",
         "My first explanation for the slowness was simply wrong, and only the stopwatch showed it. "
         "The fix I nearly shipped would have changed nothing."),
        ("Here, the explanation IS the product",
         "A right answer nobody can explain is worthless to a regulator. Showing the working is not a "
         "nice extra bolted on at the end — it is the thing being delivered."),
    ]
    for i, (h, b) in enumerate(lessons):
        card(s, Inches(7.5), Inches(2.0) + Inches(0.86 * i), Inches(5.23), Inches(0.76),
             h, b, hsize=10, bsize=8.5, hcolor=ACCENT, fill=PANEL2)

    ny = Inches(5.5)
    label(s, Inches(7.5), ny, Inches(5.23), "what is still weak → and how we would fix it")
    nexts = [
        ("Takes ~45–90s per customer", "faster hardware, a smaller model for the easy steps, and running many customers side by side — realistically under 15s"),
        ("Same input can vary slightly", "lock the model version and keep a full record of every question asked"),
        ("Finds similar cases crudely", "swap in proper meaning-based search as the case history grows"),
        ("The 60% cut-off is a guess", "tune it against real corrections once enough have piled up"),
    ]
    for i, (lim, nxt) in enumerate(nexts):
        ly = ny + Inches(0.26) + Inches(0.4 * i)
        rect(s, Inches(7.5), ly, Inches(5.23), Inches(0.36),
             fill=PANEL if i % 2 == 0 else BG, outline=None)
        text(s, Inches(7.62), ly + Inches(0.03), Inches(2.05), Inches(0.3),
             lim, size=8.5, color=MED, bold=True)
        text(s, Inches(9.75), ly + Inches(0.03), Inches(2.88), Inches(0.32),
             nxt, size=8, color=MUTED, line=1.1)

    note(s, "Close on what broke and what it taught, not on the wins. Admitting you found a silent "
            "failure — and that you were wrong once and the data corrected you — lands better than "
            "claiming it is finished.\n\n"
            "On latency, if pressed: the 45–90s is roughly 11 AI calls one after another. Most of it "
            "is waiting on a shared public API. On dedicated infrastructure — the model hosted "
            "in-house, a smaller fast model for routine steps, cached prompts, and many customers "
            "processed side by side — the same investigation lands comfortably under 15 seconds "
            "without changing any of the logic.\n\n"
            "If asked about engineering choices: no hard-coded scoring formula · watchlist screening "
            "deliberately left out of scope (the brief said “external alerts”) · exact decimal money, "
            "never floats · the test data regenerates identically every time · a permanent record of "
            "clears as well as escalations · the scratchpad is wiped on every exit path.")
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in (slide1, slide2, slide3, slide4, slide5):
        fn(prs)

    # ponytail: one assert instead of a test file — a layout regression fails the build.
    for n, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            assert sh.left >= 0 and sh.top >= 0 and sh.left + sh.width <= W \
                and sh.top + sh.height <= H, \
                f"slide {n}: shape outside canvas at ({sh.left/914400:.2f}, {sh.top/914400:.2f})"

    prs.save(OUT)
    print(f"built {len(prs.slides._sldIdLst)} slides -> {OUT}")


if __name__ == "__main__":
    main()
